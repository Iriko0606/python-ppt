# -*- coding: utf-8 -*-
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx import Presentation

def set_slide_color(slide):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(229, 255, 255)


def set_hyperlink(text_frame, hlink_text, hlink_address):
    p = text_frame.paragraphs[0]
    p.runs[0].text = hlink_text
    p.runs[0].hyperlink.address = hlink_address


def set_table_weekly_data(table):
    date = table.cell(1, 0)
    date.text = "2020/11/24"
    date.fill.solid()
    # セルの背景色変更
    date.fill.fore_color.rgb = RGBColor(255, 204, 102)
    # 輝度調整
    date.fill.fore_color.brightness = -0.25
    # theme_colorで色味調整するパターン
    # date.fill.fore_color.theme_color = MSO_THEME_COLOR.LIGHT_2
    campaign_name = table.cell(1, 1)
    campaign_name.text = "testキャンペーン"
    # フォントサイズ調整
    campaign_name.text_frame.paragraphs[0].font.size = Pt(12)
    # テーブルに実績入れる
    imp = table.cell(1, 2)
    imp.text = "10"
    # フォントの色味変える
    imp.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 153, 153)
    click = table.cell(1, 3)
    click.text = "5"


def replace_img(slide, shape):
    # テンプレートにあった画像を新しい画像に差し替え
    if shape.name == 'company_image':
        new_shape = slide.shapes.add_picture(
            "../powerpoint/image/google.png",
            Inches(shape.left.inches),
            Inches(shape.top.inches),
            Inches(shape.width.inches),
            Inches(shape.height.inches)
        )
    elif shape.name == 'media_logo':
        new_shape = slide.shapes.add_picture(
            "../powerpoint/image/googlelogo.jpg",
            Inches(shape.left.inches),
            Inches(shape.top.inches),
            Inches(shape.width.inches),
            Inches(shape.height.inches)
        )
    old_pic = shape._element
    new_pic = new_shape._element
    old_pic.addnext(new_pic)
    old_pic.getparent().remove(old_pic)


def set_contents(slide, slide_id):
    for shape in slide.shapes:
        if slide_id == 0:
            if shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER:
                # shape.nameで分岐かけるパターン
                if shape.name == 'main_title':
                    shape.text_frame.text = "CAPCAPCAP"
                    shape.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
                if shape.name == 'sub_title':
                    shape.text_frame.text = "11月週次振り返り資料作成"
                    shape.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
            elif shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
                # shape.text_frame.textで分岐かけるパターン
                if shape.text_frame.text == "{company_name}":
                    shape.text_frame.text = "CyberAgent"
                if shape.text_frame.text == "{date}":
                    shape.text_frame.text = "2020/11/20"
                if shape.text_frame.text == "{url}":
                    set_hyperlink(shape.text_frame, 'link to python-pptx @ GitHub',
                                  'https://github.com/scanny/python-pptx')
            elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                replace_img(slide, shape)
        elif slide_id == 1:
            if shape.has_table:
                if shape.name == 'weekly_performance':
                    set_table_weekly_data(shape.table)
                elif shape.name == 'monthly_performance':
                    print(shape.name)
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                replace_img(slide, shape)

def main():
    prs = Presentation("../powerpoint/base.pptx")

    page_num = len(prs.slides)
    slide_id = 0

    while slide_id <= page_num - 1:
        # スライドの枚数分ぐるぐる
        slide = prs.slides[slide_id]
        print(slide)
        # 背景色設定
        set_slide_color(slide)
        # # スライドの内容入れていく
        set_contents(slide, slide_id)
        slide_id += 1

# 出力
    prs.save("../powerpoint/test.pptx")

if __name__ == "__main__":
    main()